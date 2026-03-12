"""
File extraction service for parsing PDF, DOCX, and PPTX files.
Extracts text content from various document formats.
"""
import asyncio
from pathlib import Path
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from exceptions.custom_exceptions import FileProcessingException
from utils.file_utils import get_file_extension


async def extract_text_from_pdf(file_path: str) -> str:
    """
    Extract text content from PDF file.
    
    Args:
        file_path: Path to PDF file
        
    Returns:
        Extracted text content
        
    Raises:
        FileProcessingException: If PDF extraction fails
    """
    try:
        loop = asyncio.get_event_loop()
        text = await loop.run_in_executor(None, _extract_pdf_text, file_path)
        return text
    except Exception as e:
        raise FileProcessingException(f"Failed to extract text from PDF: {str(e)}")


def _extract_pdf_text(file_path: str) -> str:
    """Synchronous PDF text extraction."""
    text_content = []
    try:
        with open(file_path, 'rb') as file:
            pdf_reader = PdfReader(file)
            for page_num, page in enumerate(pdf_reader.pages):
                page_text = page.extract_text()
                if page_text:
                    text_content.append(page_text)
    except Exception as e:
        raise Exception(f"Error reading PDF: {str(e)}")
    
    return "\n".join(text_content)


async def extract_text_from_docx(file_path: str) -> str:
    """
    Extract text content from DOCX file.
    
    Args:
        file_path: Path to DOCX file
        
    Returns:
        Extracted text content
        
    Raises:
        FileProcessingException: If DOCX extraction fails
    """
    try:
        loop = asyncio.get_event_loop()
        text = await loop.run_in_executor(None, _extract_docx_text, file_path)
        return text
    except Exception as e:
        raise FileProcessingException(f"Failed to extract text from DOCX: {str(e)}")


def _extract_docx_text(file_path: str) -> str:
    """Synchronous DOCX text extraction."""
    text_content = []
    try:
        doc = Document(file_path)
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_content.append(paragraph.text)
        
        # Also extract text from tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        text_content.append(cell.text)
    except Exception as e:
        raise Exception(f"Error reading DOCX: {str(e)}")
    
    return "\n".join(text_content)


async def extract_text_from_pptx(file_path: str) -> str:
    """
    Extract text content from PPTX/PPT file.
    
    Args:
        file_path: Path to PPTX file
        
    Returns:
        Extracted text content
        
    Raises:
        FileProcessingException: If PPTX extraction fails
    """
    try:
        loop = asyncio.get_event_loop()
        text = await loop.run_in_executor(None, _extract_pptx_text, file_path)
        return text
    except Exception as e:
        raise FileProcessingException(f"Failed to extract text from PPTX: {str(e)}")


def _extract_pptx_text(file_path: str) -> str:
    """Synchronous PPTX text extraction."""
    text_content = []
    try:
        presentation = Presentation(file_path)
        for slide_num, slide in enumerate(presentation.slides):
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_content.append(shape.text)
    except Exception as e:
        raise Exception(f"Error reading PPTX: {str(e)}")
    
    return "\n".join(text_content)


async def extract_text_from_file(file_path: str) -> str:
    """
    Extract text from any supported file format.
    Routes to appropriate extraction function based on file extension.
    
    Args:
        file_path: Path to file
        
    Returns:
        Extracted text content
        
    Raises:
        FileProcessingException: If extraction fails or format is unsupported
    """
    file_ext = get_file_extension(file_path).lower()
    
    if file_ext == "pdf":
        return await extract_text_from_pdf(file_path)
    elif file_ext in ("docx", "doc"):
        return await extract_text_from_docx(file_path)
    elif file_ext in ("pptx", "ppt"):
        return await extract_text_from_pptx(file_path)
    else:
        raise FileProcessingException(f"Unsupported file format: {file_ext}")
